import os, subprocess, shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

print("── Starting Lightweight Multi-Format A1 Poster Generation ──")

def get_html_content(p="assets/"):
    fence_piece = f'<img class="fence-piece" src="{p}animal_pen.png" alt="Fence">' * 26
    fence_divider = f'''    <div class="fence-divider">
      <div class="fence-continuous-row">
        {fence_piece}
      </div>
    </div>'''

    return f"""<!DOCTYPE html>
<html lang="zh-TW">
<head>
  <meta charset="UTF-8" />
  <title>A1 遊戲展演海報 — 可編輯版本 (MEDIEVIL)</title>
  <link rel="preconnect" href="https://fonts.googleapis.com">
  <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
  <link href="https://fonts.googleapis.com/css2?family=Cinzel:wght@700;900&family=Noto+Serif+TC:wght@600;900&family=Plus+Jakarta+Sans:wght@500;700;800;900&display=swap" rel="stylesheet">
  <style>
    * {{
      box-sizing: border-box;
      margin: 0;
      padding: 0;
      -webkit-print-color-adjust: exact !important;
      print-color-adjust: exact !important;
      color-adjust: exact !important;
    }}
    @page {{
      size: 1360px 1923px;
      margin: 0;
    }}
    body {{
      background: #090a10;
      font-family: 'Plus Jakarta Sans', 'Noto Serif TC', -apple-system, sans-serif;
      display: flex;
      flex-direction: column;
      align-items: center;
      padding: 30px 0;
      color: #fff;
      -webkit-print-color-adjust: exact !important;
      print-color-adjust: exact !important;
    }}

    .toolbar {{
      position: fixed;
      top: 15px;
      z-index: 1000;
      background: rgba(18, 20, 32, 0.95);
      border: 1px solid rgba(245, 197, 66, 0.5);
      padding: 12px 28px;
      border-radius: 50px;
      box-shadow: 0 10px 30px rgba(0,0,0,0.85);
      display: flex;
      align-items: center;
      gap: 16px;
      backdrop-filter: blur(10px);
    }}
    .toolbar button {{
      background: linear-gradient(135deg, #f5c542, #ff7b00);
      color: #0d0e15;
      font-weight: 900;
      border: none;
      padding: 10px 22px;
      border-radius: 20px;
      cursor: pointer;
      font-size: 15px;
      box-shadow: 0 4px 15px rgba(245, 197, 66, 0.4);
    }}

    .poster-canvas {{
      width: 1360px;
      height: 1923px;
      background: #080910;
      position: relative;
      overflow: hidden;
      box-shadow: 0 30px 80px rgba(0,0,0,0.95);
      border: 5px solid #f5c542;
      border-radius: 16px;
      padding: 34px 44px;
      display: flex;
      flex-direction: column;
      justify-content: space-between;
      -webkit-print-color-adjust: exact !important;
      print-color-adjust: exact !important;
    }}

    .bg-day {{
      position: absolute;
      top: 0; left: 0; width: 100%; height: 50%;
      background: radial-gradient(circle at 20% 20%, #2b3d22 0%, #111a14 60%, #080910 100%);
      pointer-events: none;
      -webkit-print-color-adjust: exact !important;
      print-color-adjust: exact !important;
    }}
    .bg-night {{
      position: absolute;
      bottom: 0; left: 0; width: 100%; height: 58%;
      background: radial-gradient(circle at 80% 80%, #2e183a 0%, #150d22 55%, #080910 100%);
      clip-path: polygon(0 16%, 100% 0, 100% 100%, 0 100%);
      pointer-events: none;
      -webkit-print-color-adjust: exact !important;
      print-color-adjust: exact !important;
    }}
    .inner-border {{
      position: absolute;
      top: 15px; left: 15px; right: 15px; bottom: 15px;
      border: 2px solid rgba(245, 197, 66, 0.35);
      border-radius: 10px;
      pointer-events: none;
      z-index: 10;
    }}

    /* Header */
    .header-area {{
      position: relative;
      z-index: 20;
      text-align: center;
    }}
    .pill-badge {{
      display: inline-block;
      background: rgba(245, 197, 66, 0.18);
      border: 1.5px solid rgba(245, 197, 66, 0.6);
      color: #ffe885;
      padding: 6px 26px;
      border-radius: 30px;
      font-size: 16px;
      font-weight: 800;
      letter-spacing: 2px;
      margin-bottom: 4px;
      outline: none;
    }}
    .main-title {{
      font-family: 'Cinzel', serif;
      font-size: 76px;
      font-weight: 900;
      color: #f5c542;
      letter-spacing: 8px;
      text-shadow: 0 4px 30px rgba(0,0,0,0.9), 0 0 35px rgba(245, 197, 66, 0.5);
      line-height: 1.0;
      outline: none;
    }}

    /* Dynamic Lively Character Parade Under Title */
    .header-character-parade {{
      display: flex;
      justify-content: center;
      align-items: flex-end;
      gap: 18px;
      margin: 12px auto 8px;
      padding: 10px 24px;
      background: rgba(14, 18, 28, 0.75);
      border-radius: 50px;
      border: 2px solid rgba(245, 197, 66, 0.45);
      max-width: 1220px;
      backdrop-filter: blur(10px);
      box-shadow: 0 10px 30px rgba(0,0,0,0.6);
    }}
    .parade-unit {{
      display: flex;
      flex-direction: column;
      align-items: center;
      transition: transform 0.2s;
    }}
    .parade-unit img {{
      width: 56px;
      height: 56px;
      object-fit: contain;
      filter: drop-shadow(0 4px 10px rgba(0,0,0,0.85));
    }}
    .parade-name {{
      font-size: 13px;
      font-weight: 900;
      color: #e2e8f0;
      margin-top: 3px;
      outline: none;
    }}

    /* Continuous Seamless Fence Section Divider */
    .fence-divider {{
      display: flex;
      align-items: center;
      justify-content: center;
      margin: 8px auto;
      position: relative;
      z-index: 20;
      max-width: 1240px;
      overflow: hidden;
      height: 38px;
    }}
    .fence-continuous-row {{
      display: flex;
      align-items: center;
      justify-content: center;
      width: 100%;
      position: relative;
    }}
    .fence-continuous-row::before {{
      content: '';
      position: absolute;
      left: 10px; right: 10px; top: 50%;
      height: 5px;
      background: #a06e3b;
      border-radius: 3px;
      z-index: 1;
    }}
    .fence-piece {{
      width: 48px;
      height: 48px;
      object-fit: contain;
      margin: 0 -3px;
      position: relative;
      z-index: 2;
      filter: drop-shadow(0 3px 6px rgba(0,0,0,0.7));
    }}

    /* Lore Quote */
    .quote-ribbon {{
      margin: 4px auto 0;
      max-width: 1180px;
      background: rgba(10, 12, 20, 0.8);
      border-left: 5px solid #f5c542;
      border-right: 5px solid #f5c542;
      padding: 8px 22px;
      border-radius: 8px;
      font-size: 15.5px;
      line-height: 1.45;
      color: #cbd5e1;
      backdrop-filter: blur(8px);
      outline: none;
    }}

    /* Screenshots Area */
    .screenshots-area {{
      position: relative;
      z-index: 20;
      display: flex;
      justify-content: center;
      gap: 32px;
    }}
    .polaroid-card {{
      background: #fff;
      padding: 10px 10px 24px;
      border-radius: 10px;
      box-shadow: 0 18px 45px rgba(0,0,0,0.9);
      width: 560px;
    }}
    .polaroid-card.day {{ transform: rotate(-1.5deg); border: 4px solid #68d391; }}
    .polaroid-card.night {{ transform: rotate(1.5deg); border: 4px solid #e53e3e; }}
    .polaroid-img {{ width: 100%; height: 235px; object-fit: cover; border-radius: 6px; display: block; }}
    .polaroid-caption {{ color: #0f172a; font-weight: 900; font-size: 16px; margin-top: 8px; text-align: center; outline: none; }}

    /* Heroes Showcase */
    .heroes-showcase {{
      position: relative;
      z-index: 25;
      display: flex;
      justify-content: space-between;
      gap: 18px;
    }}
    .hero-sticker-card {{
      display: flex;
      align-items: center;
      gap: 15px;
      background: rgba(18, 22, 36, 0.9);
      border: 2px solid rgba(245, 197, 66, 0.4);
      padding: 12px 16px;
      border-radius: 14px;
      flex: 1;
      backdrop-filter: blur(10px);
    }}
    .hero-avatar-box {{
      width: 68px;
      height: 68px;
      background: rgba(255,255,255,0.06);
      border-radius: 12px;
      display: flex;
      align-items: center;
      justify-content: center;
      flex-shrink: 0;
      border: 1.5px solid rgba(245, 197, 66, 0.5);
    }}
    .hero-avatar-box img {{ width: 56px; height: 56px; object-fit: contain; }}
    .hero-info-title {{ font-size: 18px; font-weight: 900; color: #ffe885; outline: none; }}
    .hero-info-tag {{ font-size: 13px; color: #56ccf2; font-weight: 800; margin: 1px 0; outline: none; }}
    .hero-info-desc {{ font-size: 12.5px; color: #94a3b8; line-height: 1.35; outline: none; }}

    /* Task Priority Feature Section */
    .task-priority-showcase {{
      position: relative;
      z-index: 25;
      background: radial-gradient(circle at 20% 50%, rgba(20, 40, 30, 0.95), rgba(14, 18, 28, 0.95));
      border: 2px solid #2ec4b6;
      border-radius: 16px;
      padding: 12px 20px;
      box-shadow: 0 8px 25px rgba(46, 196, 182, 0.25);
    }}
    .task-head {{
      display: flex;
      justify-content: space-between;
      align-items: center;
      margin-bottom: 8px;
    }}
    .task-title {{
      font-size: 19px;
      font-weight: 900;
      color: #a7f3d0;
      display: flex;
      align-items: center;
      gap: 8px;
      outline: none;
    }}
    .task-badge {{
      background: rgba(46, 196, 182, 0.2);
      border: 1px solid #2ec4b6;
      color: #a7f3d0;
      font-size: 12.5px;
      font-weight: 800;
      padding: 3px 14px;
      border-radius: 20px;
      outline: none;
    }}
    .task-grid {{
      display: grid;
      grid-template-columns: repeat(3, 1fr);
      gap: 14px;
    }}
    .task-card {{
      background: rgba(10, 12, 20, 0.6);
      border-radius: 12px;
      padding: 9px 12px;
      border: 1.5px solid rgba(46, 196, 182, 0.35);
      display: flex;
      align-items: center;
      gap: 12px;
    }}
    .task-icon-box {{
      width: 48px;
      height: 48px;
      background: rgba(46, 196, 182, 0.15);
      border: 1px solid #2ec4b6;
      border-radius: 10px;
      display: flex;
      align-items: center;
      justify-content: center;
      flex-shrink: 0;
    }}
    .task-icon-box img {{ width: 38px; height: 38px; object-fit: contain; }}
    .task-name {{ font-size: 15px; font-weight: 900; color: #fff; outline: none; }}
    .task-desc {{ font-size: 11.5px; color: #94a3b8; line-height: 1.35; outline: none; }}

    /* Magic Strike & Skill Tree Section */
    .magic-strike-showcase {{
      position: relative;
      z-index: 25;
      background: radial-gradient(circle at 50% 50%, rgba(38, 18, 58, 0.95), rgba(14, 18, 28, 0.95));
      border: 2px solid #a855f7;
      border-radius: 16px;
      padding: 12px 20px;
      box-shadow: 0 8px 30px rgba(168, 85, 247, 0.25);
    }}
    .magic-head {{
      display: flex;
      justify-content: space-between;
      align-items: center;
      margin-bottom: 8px;
    }}
    .magic-title {{
      font-size: 19px;
      font-weight: 900;
      color: #f3e8ff;
      display: flex;
      align-items: center;
      gap: 8px;
      outline: none;
    }}
    .magic-badge {{
      background: rgba(245, 197, 66, 0.2);
      border: 1px solid #f5c542;
      color: #ffe885;
      font-size: 12.5px;
      font-weight: 800;
      padding: 3px 14px;
      border-radius: 20px;
      outline: none;
    }}
    .magic-grid {{
      display: grid;
      grid-template-columns: repeat(3, 1fr);
      gap: 14px;
    }}
    .magic-card {{
      background: rgba(10, 12, 20, 0.6);
      border-radius: 12px;
      padding: 9px 12px;
      border: 1.5px solid rgba(255,255,255,0.1);
      display: flex;
      align-items: center;
      gap: 12px;
    }}
    .skill-icon-box {{
      width: 48px;
      height: 48px;
      border-radius: 10px;
      display: flex;
      align-items: center;
      justify-content: center;
      flex-shrink: 0;
      border: 1.5px solid rgba(255,255,255,0.2);
    }}
    .skill-icon-box img {{ width: 38px; height: 38px; object-fit: contain; }}
    .magic-card.fire {{ border-color: rgba(249, 115, 22, 0.5); }}
    .magic-card.fire .skill-icon-box {{ background: rgba(249, 115, 22, 0.2); border-color: #f97316; }}
    .magic-card.lightning {{ border-color: rgba(234, 179, 8, 0.5); }}
    .magic-card.lightning .skill-icon-box {{ background: rgba(234, 179, 8, 0.2); border-color: #eab308; }}
    .magic-card.freeze {{ border-color: rgba(56, 189, 248, 0.5); }}
    .magic-card.freeze .skill-icon-box {{ background: rgba(56, 189, 248, 0.2); border-color: #38bdf8; }}
    .magic-name {{ font-size: 15px; font-weight: 900; color: #fff; outline: none; }}
    .magic-desc {{ font-size: 11.5px; color: #94a3b8; line-height: 1.35; outline: none; }}

    /* Split Features */
    .split-features {{
      position: relative;
      z-index: 20;
      display: grid;
      grid-template-columns: 1fr 1fr;
      gap: 20px;
    }}
    .feature-panel {{
      background: rgba(14, 17, 28, 0.9);
      border-radius: 14px;
      padding: 12px 18px;
      border: 1.5px solid rgba(255,255,255,0.12);
    }}
    .panel-night {{ border-top: 4px solid #e63946; }}
    .panel-day {{ border-top: 4px solid #2ec4b6; }}
    .panel-head {{ font-size: 17px; font-weight: 900; margin-bottom: 4px; color: #fff; outline: none; }}
    .panel-desc {{ font-size: 12.5px; color: #cbd5e1; line-height: 1.4; margin-bottom: 8px; outline: none; }}

    .sprite-ribbon {{
      display: flex;
      justify-content: space-between;
      gap: 8px;
      background: rgba(0,0,0,0.35);
      padding: 6px 8px;
      border-radius: 8px;
    }}
    .sprite-mini-card {{ text-align: center; flex: 1; }}
    .sprite-mini-card img {{ width: 40px; height: 40px; object-fit: contain; display: block; margin: 0 auto 2px; }}
    .sprite-mini-card span {{ font-size: 11.5px; font-weight: 800; color: #ffe885; display: block; outline: none; }}

    /* Bottom Biomes */
    .bottom-showcase {{
      position: relative;
      z-index: 20;
      background: rgba(14, 17, 28, 0.9);
      border-radius: 14px;
      padding: 12px 18px;
      border: 1.5px solid rgba(245, 197, 66, 0.4);
    }}
    .biomes-grid {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 12px; }}
    .biome-item {{ display: flex; align-items: center; gap: 8px; background: rgba(255,255,255,0.05); padding: 6px 10px; border-radius: 8px; }}
    .biome-item img {{ width: 38px; height: 38px; object-fit: contain; }}
    .biome-name {{ font-size: 13.5px; font-weight: 900; color: #ffe885; outline: none; }}
    .biome-trait {{ font-size: 11px; color: #94a3b8; outline: none; }}

    .poster-footer {{
      margin-top: 8px;
      text-align: center;
      position: relative;
      z-index: 20;
      font-size: 14px;
      font-weight: 900;
      color: #f5c542;
      letter-spacing: 2px;
      outline: none;
    }}
    @media print {{
      html, body {{
        width: 1360px !important;
        height: 1923px !important;
        margin: 0 !important;
        padding: 0 !important;
        background: #080910 !important;
        -webkit-print-color-adjust: exact !important;
        print-color-adjust: exact !important;
      }}
      .toolbar {{ display: none !important; }}
      .poster-canvas {{
        width: 1360px !important;
        height: 1923px !important;
        margin: 0 !important;
        padding: 34px 44px !important;
        border: 5px solid #f5c542 !important;
        border-radius: 0 !important;
        box-shadow: none !important;
        page-break-inside: avoid !important;
        break-inside: avoid !important;
        -webkit-print-color-adjust: exact !important;
        print-color-adjust: exact !important;
      }}
    }}
  </style>
</head>
<body>

  <div class="toolbar">
    <span style="font-size:14px; color:#cbd5e1;">💡 提示：點擊海報中任何文字即可<b>直接修改編輯</b></span>
    <button onclick="window.print()">🖨️ 列印 / 另存 A1 PDF</button>
  </div>

  <div class="poster-canvas">
    <div class="bg-day"></div>
    <div class="bg-night"></div>
    <div class="inner-border"></div>

    <div class="header-area">
      <div class="pill-badge" contenteditable="true">✦ 2026 年度暗黑奇幻像素殖民地生存策略 ✦</div>
      <h1 class="main-title" contenteditable="true">M E D I E V I L</h1>

      <div class="header-character-parade">
        <div class="parade-unit" style="transform: rotate(-6deg);">
          <img src="{p}farmer.png" alt="Farmer">
          <span class="parade-name" contenteditable="true">拓荒農夫</span>
        </div>
        <div class="parade-unit" style="transform: translateY(-8px);">
          <img src="{p}horse.png" alt="Horse">
          <span class="parade-name" style="color:#56ccf2;" contenteditable="true">疾行駿馬</span>
        </div>
        <div class="parade-unit" style="transform: rotate(5deg);">
          <img src="{p}boar.png" alt="Boar">
          <span class="parade-name" style="color:#68d391;" contenteditable="true">野豬產肉</span>
        </div>
        <div class="parade-unit" style="transform: scale(1.18) translateY(-10px);">
          <img src="{p}knight.png" alt="Knight">
          <span class="parade-name" style="color:#ffe885;" contenteditable="true">聖殿騎士</span>
        </div>
        <div class="parade-unit" style="transform: scale(1.15) translateY(-8px);">
          <img src="{p}magician.png" alt="Magician">
          <span class="parade-name" style="color:#c084fc;" contenteditable="true">元素法師</span>
        </div>
        <div class="parade-unit" style="transform: rotate(-4deg);">
          <img src="{p}flying_squirrel.png" alt="Squirrel">
          <span class="parade-name" contenteditable="true">隨行飛鼠</span>
        </div>
        <div class="parade-unit" style="transform: rotate(-5deg) translateY(-6px);">
          <img src="{p}werewolf.png" alt="Werewolf">
          <span class="parade-name" style="color:#e53e3e;" contenteditable="true">狂暴狼人</span>
        </div>
        <div class="parade-unit" style="transform: rotate(5deg) translateY(-8px);">
          <img src="{p}vampire.png" alt="Vampire">
          <span class="parade-name" style="color:#a855f7;" contenteditable="true">暗夜吸血鬼</span>
        </div>
        <div class="parade-unit" style="transform: rotate(-3deg);">
          <img src="{p}zombie.png" alt="Zombie">
          <span class="parade-name" style="color:#a0aec0;" contenteditable="true">劇毒殭屍</span>
        </div>
        <div class="parade-unit" style="transform: rotate(4deg) translateY(-4px);">
          <img src="{p}bear.png" alt="Bear">
          <span class="parade-name" style="color:#f6ad55;" contenteditable="true">狂暴巨熊</span>
        </div>
      </div>

      <div class="quote-ribbon" contenteditable="true">
        「白晝在未知的迷霧荒野中拓荒建設、採集資源、耕種農田、馴服野獸；黑夜深淵魔物破曉襲擊之時，詠唱全圖魔法抵禦狂暴狼人、吸血鬼與劇毒殭屍的無盡圍攻。<strong>白晝拓荒寸土，黑夜寸步不讓！</strong>」
      </div>
    </div>

    {fence_divider}

    <div class="screenshots-area">
      <div class="polaroid-card day">
        <img class="polaroid-img" src="{p}screenshot_day.png" alt="Day Gameplay">
        <div class="polaroid-caption" contenteditable="true">☀️ 白晝營地：農耕建設與生態開拓 (120s)</div>
      </div>
      <div class="polaroid-card night">
        <img class="polaroid-img" src="{p}screenshot_night.png" alt="Night Gameplay">
        <div class="polaroid-caption" contenteditable="true">🌙 暗夜防線：魔物夜襲與全圖魔法交鋒 (60s)</div>
      </div>
    </div>

    {fence_divider}

    <div class="heroes-showcase">
      <div class="hero-sticker-card">
        <div class="hero-avatar-box"><img src="{p}farmer.png" alt="Farmer"></div>
        <div>
          <div class="hero-info-title" contenteditable="true">拓荒農夫</div>
          <div class="hero-info-tag" contenteditable="true">經濟核心 · 馴化專家</div>
          <div class="hero-info-desc" contenteditable="true">工作神速、小麥耕作採收，享有 1.5 倍野生動物馴服成功率！</div>
        </div>
      </div>

      <div class="hero-sticker-card">
        <div class="hero-avatar-box"><img src="{p}knight.png" alt="Knight"></div>
        <div>
          <div class="hero-info-title" contenteditable="true">聖殿騎士</div>
          <div class="hero-info-tag" contenteditable="true">前線護衛 · 狩獵大師</div>
          <div class="hero-info-desc" contenteditable="true">身披重甲堅盾在前線抗怪，獵捕野獸享有 50% 爆擊與高額傷害！</div>
        </div>
      </div>

      <div class="hero-sticker-card">
        <div class="hero-avatar-box"><img src="{p}magician.png" alt="Magician"></div>
        <div>
          <div class="hero-info-title" contenteditable="true">元素法師</div>
          <div class="hero-info-tag" contenteditable="true">戰略砲台 · 元素掌控</div>
          <div class="hero-info-desc" contenteditable="true">詠唱火焰術範圍爆破、連鎖閃電穿透與冰凍霜結，夜間群攻核心！</div>
        </div>
      </div>
    </div>

    {fence_divider}

    <div class="task-priority-showcase">
      <div class="task-head">
        <div class="task-title" contenteditable="true">
          <span>📋⚡【高度自由的任務安排優先順序，由你決定！！】</span>
        </div>
        <div class="task-badge" contenteditable="true">✦ 智慧排程佇列 · 隨心調度全殖民地 ✦</div>
      </div>
      <div class="task-grid">
        <div class="task-card">
          <div class="task-icon-box"><img src="{p}crop.png" alt="Crop" /></div>
          <div>
            <div class="task-name" contenteditable="true">🌾 農耕優先 vs 🪵 伐木採石</div>
            <div class="task-desc" contenteditable="true">自由分配居民作業優先級，糧食儲備與建材擴張完全掌控！</div>
          </div>
        </div>
        <div class="task-card">
          <div class="task-icon-box"><img src="{p}meat.png" alt="Meat" /></div>
          <div>
            <div class="task-name" contenteditable="true">🐗 野外狩獵 vs 🐎 馴化畜牧</div>
            <div class="task-desc" contenteditable="true">優先獵肉充飢或圈養駿馬全域加速？策略節奏由你抉擇！</div>
          </div>
        </div>
        <div class="task-card">
          <div class="task-icon-box"><img src="{p}tower.png" alt="Tower" /></div>
          <div>
            <div class="task-name" contenteditable="true">🏰 要塞修築 vs 🛡️ 前線巡邏</div>
            <div class="task-desc" contenteditable="true">白晝全力趕工城牆箭塔，入夜重裝騎士前線阻擊夜襲！</div>
          </div>
        </div>
      </div>
    </div>

    {fence_divider}

    <div class="magic-strike-showcase">
      <div class="magic-head">
        <div class="magic-title" contenteditable="true">
          <span>⚡🔥❄️ 全圖三大元素法術天罰打擊 (Global Magic Strike)</span>
        </div>
        <div class="magic-badge" contenteditable="true">✦ 即時詠唱 · 全圖無死角精準轟炸 ✦</div>
      </div>
      <div class="magic-grid">
        <div class="magic-card fire">
          <div class="skill-icon-box"><img src="{p}scorched.png" alt="Fire" /></div>
          <div>
            <div class="magic-name" contenteditable="true">火焰天罰 (Meteor Fire)</div>
            <div class="magic-desc" contenteditable="true">全圖指定區域隕石轟炸，高額範圍爆破與持續熔岩灼燒！</div>
          </div>
        </div>
        <div class="magic-card lightning">
          <div class="skill-icon-box"><img src="{p}magician.png" alt="Lightning" /></div>
          <div>
            <div class="magic-name" contenteditable="true">連鎖閃電 (Chain Lightning)</div>
            <div class="magic-desc" contenteditable="true">九天神雷貫穿魔物群，多重連鎖彈射瞬秒高速突襲夜行者！</div>
          </div>
        </div>
        <div class="magic-card freeze">
          <div class="skill-icon-box"><img src="{p}swamp.png" alt="Freeze" /></div>
          <div>
            <div class="magic-name" contenteditable="true">極地冰凍 (Deep Freeze)</div>
            <div class="magic-desc" contenteditable="true">3x3 暴風雪大範圍霜結凍結，強制冰封所有狂暴魔物！</div>
          </div>
        </div>
      </div>

      <div style="margin-top: 10px; padding-top: 8px; border-top: 1px dashed rgba(245, 197, 66, 0.35); display: flex; justify-content: space-between; align-items: center;">
        <div class="magic-title" style="color:#ffe885;" contenteditable="true">
          <span>🌟 獨特技能點天賦升級系統 (Skill Tree Progression)</span>
        </div>
        <div class="magic-badge" style="background:rgba(86, 204, 242, 0.2); border-color:#56ccf2; color:#56ccf2;" contenteditable="true">✦ 按 [K] 鍵開啟天賦樹 · 每夜黎明結算 2 SP 點數 ✦</div>
      </div>
      <div class="magic-grid" style="margin-top: 8px;">
        <div class="magic-card" style="border-color: rgba(104, 211, 145, 0.5);">
          <div class="skill-icon-box" style="background:rgba(104, 211, 145, 0.15); border-color:#68d391;"><img src="{p}farmer.png" alt="Economy" /></div>
          <div>
            <div class="magic-name" style="color:#a7f3d0;" contenteditable="true">🌾 生產天賦 · 馴獸大師</div>
            <div class="magic-desc" contenteditable="true">作物生長速度 +25%，採集雙倍資源，野生動物馴化成功率大增！</div>
          </div>
        </div>
        <div class="magic-card" style="border-color: rgba(245, 197, 66, 0.5);">
          <div class="skill-icon-box" style="background:rgba(245, 197, 66, 0.15); border-color:#f5c542;"><img src="{p}knight.png" alt="Defense" /></div>
          <div>
            <div class="magic-name" style="color:#fef08a;" contenteditable="true">🛡️ 防禦天賦 · 堅甲要塞</div>
            <div class="magic-desc" contenteditable="true">防禦箭塔攻速射程 +30%，城牆耐久度提升，全員拓荒者生命上限強化！</div>
          </div>
        </div>
        <div class="magic-card" style="border-color: rgba(192, 132, 252, 0.5);">
          <div class="skill-icon-box" style="background:rgba(192, 132, 252, 0.15); border-color:#c084fc;"><img src="{p}raw_stone.png" alt="Magic" /></div>
          <div>
            <div class="magic-name" style="color:#f3e8ff;" contenteditable="true">🔮 魔法天賦 · 元素超載</div>
            <div class="magic-desc" contenteditable="true">全圖法術冷卻縮短 25%，火焰範圍爆破擴大 50%，冰凍控制時間延長！</div>
          </div>
        </div>
      </div>
    </div>

    {fence_divider}

    <div class="split-features">
      <div class="feature-panel panel-night">
        <div class="panel-head" contenteditable="true">🌙 暗夜深淵魔物狂潮與防禦箭塔</div>
        <div class="panel-desc" contenteditable="true">
          黑夜降臨巢穴將湧出嗜血魔物！防禦箭塔全自動遠程齊射，配合全圖元素法術大範圍殲滅敵軍！
        </div>
        <div class="sprite-ribbon">
          <div class="sprite-mini-card">
            <img src="{p}werewolf.png" alt="Werewolf">
            <span contenteditable="true">狂暴狼人</span>
          </div>
          <div class="sprite-mini-card">
            <img src="{p}vampire.png" alt="Vampire">
            <span contenteditable="true">暗夜吸血鬼</span>
          </div>
          <div class="sprite-mini-card">
            <img src="{p}zombie.png" alt="Zombie">
            <span contenteditable="true">劇毒殭屍</span>
          </div>
          <div class="sprite-mini-card">
            <img src="{p}tower.png" alt="Tower">
            <span contenteditable="true">防禦箭塔</span>
          </div>
        </div>
      </div>

      <div class="feature-panel panel-day">
        <div class="panel-head" contenteditable="true">🌱 白晝荒野馴化與畜欄增益經濟</div>
        <div class="panel-desc" contenteditable="true">
          探索邊境並馴服野生動物！圈養動物為營地持續產肉，圈養駿馬賦予全體居民 +37.5% 移動加速！
        </div>
        <div class="sprite-ribbon">
          <div class="sprite-mini-card">
            <img src="{p}boar.png" alt="Boar">
            <span contenteditable="true">野豬(產肉)</span>
          </div>
          <div class="sprite-mini-card">
            <img src="{p}horse.png" alt="Horse">
            <span contenteditable="true">駿馬(加速)</span>
          </div>
          <div class="sprite-mini-card">
            <img src="{p}flying_squirrel.png" alt="Squirrel">
            <span contenteditable="true">飛鼠(寵物)</span>
          </div>
          <div class="sprite-mini-card">
            <img src="{p}animal_pen.png" alt="Pen">
            <span contenteditable="true">畜欄建築</span>
          </div>
        </div>
      </div>
    </div>

    {fence_divider}

    <div class="bottom-showcase">
      <div class="biomes-grid">
        <div class="biome-item">
          <img src="{p}river.png" alt="River">
          <div>
            <div class="biome-name" contenteditable="true">水彩河流</div>
            <div class="biome-trait" contenteditable="true">自然緩速 50%</div>
          </div>
        </div>
        <div class="biome-item">
          <img src="{p}mountain.png" alt="Mountain">
          <div>
            <div class="biome-name" contenteditable="true">連綿山脈</div>
            <div class="biome-trait" contenteditable="true">泛洪全揭露障礙</div>
          </div>
        </div>
        <div class="biome-item">
          <img src="{p}swamp.png" alt="Swamp">
          <div>
            <div class="biome-name" contenteditable="true">幽暗泥沼</div>
            <div class="biome-trait" contenteditable="true">5.0s 受困倒數</div>
          </div>
        </div>
        <div class="biome-item">
          <img src="{p}scorched.png" alt="Scorched">
          <div>
            <div class="biome-name" contenteditable="true">熔岩焦土</div>
            <div class="biome-trait" contenteditable="true">持續高溫灼燒</div>
          </div>
        </div>
      </div>
    </div>

    <div class="poster-footer" contenteditable="true">
      第二組 出品 · 《MEDIEVIL》
    </div>
  </div>
</body>
</html>"""

# Write HTMLs
with open('docs/a1_poster_editable.html', 'w', encoding='utf-8') as f:
    f.write(get_html_content('assets/'))
with open('a1_poster_editable.html', 'w', encoding='utf-8') as f:
    f.write(get_html_content('docs/assets/'))
print("[OK] Successfully generated lightweight a1_poster_editable.html (~25 KB)")

# ── 2. HIGH-RES VECTOR PDF (Exact Match with HTML) ───────────────────────────
edge_exe = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
if not os.path.exists(edge_exe):
    edge_exe = r"C:\Program Files\Google\Chrome\Application\chrome.exe"

html_abs = os.path.abspath("docs/a1_poster_editable.html")
pdf_abs = os.path.abspath("a1_poster_editable.pdf")
pdf_docs_abs = os.path.abspath("docs/a1_poster_editable.pdf")

cmd = f'"{edge_exe}" --headless --disable-gpu --print-to-pdf="{pdf_abs}" --no-pdf-header-footer "file:///{html_abs}"'
res = subprocess.run(cmd, shell=True, capture_output=True)
if os.path.exists(pdf_abs):
    shutil.copy(pdf_abs, pdf_docs_abs)
    print(f"[OK] Successfully generated a1_poster_editable.pdf ({os.path.getsize(pdf_abs)//1024} KB)")

# ── 3. CLEAN VECTOR SVG ──────────────────────────────────────────────────────
def get_svg_content(p="assets/"):
    fence_svg_row = "".join([f'<image href="{p}animal_pen.png" xlink:href="{p}animal_pen.png" x="{120 + i * 90}" y="0" width="98" height="98" />' for i in range(24)])
    return f"""<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink" viewBox="0 0 2480 3508" width="2480" height="3508">
  <defs>
    <linearGradient id="dayGrad" x1="0%" y1="0%" x2="100%" y2="100%">
      <stop offset="0%" stop-color="#2a3d24" />
      <stop offset="60%" stop-color="#121c16" />
      <stop offset="100%" stop-color="#080910" />
    </linearGradient>
    <linearGradient id="nightGrad" x1="0%" y1="0%" x2="100%" y2="100%">
      <stop offset="0%" stop-color="#2d1838" />
      <stop offset="50%" stop-color="#150d22" />
      <stop offset="100%" stop-color="#080910" />
    </linearGradient>
  </defs>

  <rect width="2480" height="3508" fill="#080910" />
  <rect width="2480" height="1800" fill="url(#dayGrad)" />
  <polygon points="0,1700 2480,1200 2480,3508 0,3508" fill="url(#nightGrad)" />

  <rect x="50" y="50" width="2380" height="3408" rx="28" fill="none" stroke="#f5c542" stroke-width="6" />
  <rect x="70" y="70" width="2340" height="3368" rx="20" fill="none" stroke="#f5c542" stroke-width="2" stroke-opacity="0.4" />

  <g id="header" transform="translate(1240, 130)" text-anchor="middle">
    <rect x="-420" y="0" width="840" height="60" rx="30" fill="#1e1a2d" stroke="#f5c542" stroke-width="2" />
    <text x="0" y="42" font-family="sans-serif" font-size="32" font-weight="bold" fill="#ffe885">✦ 2026 年度暗黑奇幻像素殖民地生存策略 ✦</text>
    <text x="0" y="160" font-family="serif" font-size="120" font-weight="900" fill="#f5c542">M E D I E V I L</text>

    <g transform="translate(-1080, 190)">
      <rect x="0" y="0" width="2160" height="160" rx="45" fill="#0e121c" fill-opacity="0.85" stroke="#f5c542" stroke-width="2" />
      <g transform="translate(70, 20)">
        <image href="{p}farmer.png" xlink:href="{p}farmer.png" x="25" y="0" width="95" height="95" />
        <text x="72" y="125" text-anchor="middle" font-family="sans-serif" font-size="24" font-weight="bold" fill="#e2e8f0">拓荒農夫</text>
      </g>
      <g transform="translate(275, 15)">
        <image href="{p}horse.png" xlink:href="{p}horse.png" x="25" y="0" width="95" height="95" />
        <text x="72" y="125" text-anchor="middle" font-family="sans-serif" font-size="24" font-weight="bold" fill="#56ccf2">疾行駿馬</text>
      </g>
      <g transform="translate(480, 20)">
        <image href="{p}boar.png" xlink:href="{p}boar.png" x="25" y="0" width="95" height="95" />
        <text x="72" y="125" text-anchor="middle" font-family="sans-serif" font-size="24" font-weight="bold" fill="#68d391">野豬產肉</text>
      </g>
      <g transform="translate(695, 8)">
        <image href="{p}knight.png" xlink:href="{p}knight.png" x="25" y="0" width="115" height="115" />
        <text x="82" y="135" text-anchor="middle" font-family="sans-serif" font-size="26" font-weight="bold" fill="#ffe885">聖殿騎士</text>
      </g>
      <g transform="translate(915, 8)">
        <image href="{p}magician.png" xlink:href="{p}magician.png" x="25" y="0" width="115" height="115" />
        <text x="82" y="135" text-anchor="middle" font-family="sans-serif" font-size="26" font-weight="bold" fill="#c084fc">元素法師</text>
      </g>
      <g transform="translate(1135, 20)">
        <image href="{p}flying_squirrel.png" xlink:href="{p}flying_squirrel.png" x="25" y="0" width="95" height="95" />
        <text x="72" y="125" text-anchor="middle" font-family="sans-serif" font-size="24" font-weight="bold" fill="#e2e8f0">隨行飛鼠</text>
      </g>
      <g transform="translate(1340, 15)">
        <image href="{p}werewolf.png" xlink:href="{p}werewolf.png" x="25" y="0" width="95" height="95" />
        <text x="72" y="125" text-anchor="middle" font-family="sans-serif" font-size="24" font-weight="bold" fill="#e53e3e">狂暴狼人</text>
      </g>
      <g transform="translate(1545, 15)">
        <image href="{p}vampire.png" xlink:href="{p}vampire.png" x="25" y="0" width="95" height="95" />
        <text x="72" y="125" text-anchor="middle" font-family="sans-serif" font-size="24" font-weight="bold" fill="#a855f7">暗夜吸血鬼</text>
      </g>
      <g transform="translate(1750, 20)">
        <image href="{p}zombie.png" xlink:href="{p}zombie.png" x="25" y="0" width="95" height="95" />
        <text x="72" y="125" text-anchor="middle" font-family="sans-serif" font-size="24" font-weight="bold" fill="#a0aec0">劇毒殭屍</text>
      </g>
      <g transform="translate(1955, 15)">
        <image href="{p}bear.png" xlink:href="{p}bear.png" x="25" y="0" width="95" height="95" />
        <text x="72" y="125" text-anchor="middle" font-family="sans-serif" font-size="24" font-weight="bold" fill="#f6ad55">狂暴巨熊</text>
      </g>
    </g>

    <rect x="-1050" y="375" width="2100" height="95" rx="14" fill="#0a0c14" fill-opacity="0.85" stroke="#f5c542" stroke-width="2" />
    <text x="0" y="415" font-family="sans-serif" font-size="28" fill="#cbd5e1">「白晝在未知的迷霧荒野中拓荒建設、採集資源、耕種農田、馴服野獸；</text>
    <text x="0" y="452" font-family="sans-serif" font-size="28" fill="#ffe885" font-weight="bold">黑夜深淵魔物破曉襲擊之時，詠唱全圖魔法抵禦狂暴狼人、吸血鬼與劇毒殭屍的無盡圍攻。白晝拓荒寸土，黑夜寸步不讓！」</text>
  </g>

  <g id="fence1" transform="translate(100, 620)">
    <rect x="20" y="45" width="2240" height="8" rx="4" fill="#a06e3b" />
    {fence_svg_row}
  </g>

  <g id="screenshots" transform="translate(0, 725)">
    <g transform="translate(160, 0)">
      <rect x="0" y="0" width="1040" height="560" rx="14" fill="#ffffff" stroke="#68d391" stroke-width="5" />
      <image href="{p}screenshot_day.png" xlink:href="{p}screenshot_day.png" x="20" y="20" width="1000" height="460" preserveAspectRatio="xMidYMid slice" />
      <text x="520" y="525" text-anchor="middle" font-family="sans-serif" font-size="30" font-weight="bold" fill="#1a202c">☀️ 白晝營地：農耕建設與生態開拓 (120s)</text>
    </g>
    <g transform="translate(1280, 0)">
      <rect x="0" y="0" width="1040" height="560" rx="14" fill="#ffffff" stroke="#e53e3e" stroke-width="5" />
      <image href="{p}screenshot_night.png" xlink:href="{p}screenshot_night.png" x="20" y="20" width="1000" height="460" preserveAspectRatio="xMidYMid slice" />
      <text x="520" y="525" text-anchor="middle" font-family="sans-serif" font-size="30" font-weight="bold" fill="#1a202c">🌙 暗夜防線：魔物夜襲與全圖魔法交鋒 (60s)</text>
    </g>
  </g>

  <g id="fence2" transform="translate(100, 1310)">
    <rect x="20" y="45" width="2240" height="8" rx="4" fill="#a06e3b" />
    {fence_svg_row}
  </g>

  <g id="heroes" transform="translate(0, 1410)">
    <g transform="translate(160, 0)">
      <rect width="680" height="165" rx="20" fill="#121624" stroke="#2ec4b6" stroke-width="2.5" />
      <rect x="20" y="20" width="125" height="125" rx="16" fill="#1e2538" stroke="#f5c542" stroke-width="1.5" />
      <image href="{p}farmer.png" xlink:href="{p}farmer.png" x="33" y="33" width="100" height="100" />
      <text x="165" y="60" font-family="serif" font-size="32" font-weight="bold" fill="#ffffff">拓荒農夫 (Farmer)</text>
      <text x="165" y="95" font-family="sans-serif" font-size="23" font-weight="bold" fill="#2ec4b6">經濟核心 · 0.6x 工作神速 · 1.5x 馴服</text>
      <text x="165" y="130" font-family="sans-serif" font-size="20" fill="#cbd5e1">負責小麥耕作與生態馴化，殖民地擴張的基石！</text>
    </g>
    <g transform="translate(900, 0)">
      <rect width="680" height="165" rx="20" fill="#121624" stroke="#f5c542" stroke-width="2.5" />
      <rect x="20" y="20" width="125" height="125" rx="16" fill="#1e2538" stroke="#f5c542" stroke-width="1.5" />
      <image href="{p}knight.png" xlink:href="{p}knight.png" x="33" y="33" width="100" height="100" />
      <text x="165" y="60" font-family="serif" font-size="32" font-weight="bold" fill="#ffffff">聖殿騎士 (Knight)</text>
      <text x="165" y="95" font-family="sans-serif" font-size="23" font-weight="bold" fill="#f5c542">前線堅守 · 狩獵 50% 爆擊 · 銅牆鐵壁</text>
      <text x="165" y="130" font-family="sans-serif" font-size="20" fill="#cbd5e1">身披堅甲在前線格擋，阻擋夜間狼人與吸血鬼推進！</text>
    </g>
    <g transform="translate(1640, 0)">
      <rect width="680" height="165" rx="20" fill="#121624" stroke="#a855f7" stroke-width="2.5" />
      <rect x="20" y="20" width="125" height="125" rx="16" fill="#1e2538" stroke="#f5c542" stroke-width="1.5" />
      <image href="{p}magician.png" xlink:href="{p}magician.png" x="33" y="33" width="100" height="100" />
      <text x="165" y="60" font-family="serif" font-size="32" font-weight="bold" fill="#ffffff">元素法師 (Magician)</text>
      <text x="165" y="95" font-family="sans-serif" font-size="23" font-weight="bold" fill="#a855f7">戰略砲台 · 火焰 / 閃電 / 冰凍三大法術</text>
      <text x="165" y="130" font-family="sans-serif" font-size="20" fill="#cbd5e1">詠唱火焰範圍爆破與冰凍霜結，夜間群攻核心！</text>
    </g>
  </g>

  <g id="fence3" transform="translate(100, 1600)">
    <rect x="20" y="45" width="2240" height="8" rx="4" fill="#a06e3b" />
    {fence_svg_row}
  </g>

  <g id="task_priority" transform="translate(160, 1700)">
    <rect width="2160" height="175" rx="20" fill="#102520" stroke="#2ec4b6" stroke-width="2.5" />
    <text x="40" y="45" font-family="serif" font-size="30" font-weight="bold" fill="#a7f3d0">📋⚡【高度自由的任務安排優先順序，由你決定！！】</text>
    <rect x="1600" y="16" width="510" height="40" rx="20" fill="#1b3d35" stroke="#2ec4b6" stroke-width="1.5" />
    <text x="1855" y="43" text-anchor="middle" font-family="sans-serif" font-size="19" font-weight="bold" fill="#a7f3d0">✦ 智慧排程佇列 · 隨心調度全殖民地 ✦</text>

    <g transform="translate(40, 65)">
      <rect width="680" height="90" rx="12" fill="#0d1416" stroke="#2ec4b6" stroke-width="1.5" />
      <image href="{p}crop.png" xlink:href="{p}crop.png" x="20" y="15" width="60" height="60" />
      <text x="95" y="40" font-family="serif" font-size="22" font-weight="bold" fill="#ffffff">🌾 農耕優先 vs 🪵 伐木採石</text>
      <text x="95" y="68" font-family="sans-serif" font-size="17" fill="#94a3b8">自由分配居民作業優先級，掌控糧食與建材！</text>
    </g>
    <g transform="translate(740, 65)">
      <rect width="680" height="90" rx="12" fill="#0d1416" stroke="#2ec4b6" stroke-width="1.5" />
      <image href="{p}meat.png" xlink:href="{p}meat.png" x="20" y="15" width="60" height="60" />
      <text x="95" y="40" font-family="serif" font-size="22" font-weight="bold" fill="#ffffff">🐗 野外狩獵 vs 🐎 馴化畜牧</text>
      <text x="95" y="68" font-family="sans-serif" font-size="17" fill="#94a3b8">優先獵肉充飢或圈養駿馬全域加速？策略由你抉擇！</text>
    </g>
    <g transform="translate(1440, 65)">
      <rect width="680" height="90" rx="12" fill="#0d1416" stroke="#2ec4b6" stroke-width="1.5" />
      <image href="{p}tower.png" xlink:href="{p}tower.png" x="20" y="15" width="60" height="60" />
      <text x="95" y="40" font-family="serif" font-size="22" font-weight="bold" fill="#ffffff">🏰 要塞修築 vs 🛡️ 前線巡邏</text>
      <text x="95" y="68" font-family="sans-serif" font-size="17" fill="#94a3b8">白晝全力搶修城牆箭塔，入夜重裝騎士阻擊夜襲！</text>
    </g>
  </g>

  <g id="fence4" transform="translate(100, 1895)">
    <rect x="20" y="45" width="2240" height="8" rx="4" fill="#a06e3b" />
    {fence_svg_row}
  </g>

  <g id="magic_and_skills" transform="translate(160, 1995)">
    <rect width="2160" height="300" rx="20" fill="#1a1128" stroke="#a855f7" stroke-width="2.5" />
    <text x="40" y="40" font-family="serif" font-size="28" font-weight="bold" fill="#f3e8ff">⚡🔥❄️ 全圖三大元素法術天罰打擊 (Global Magic Strike)</text>
    <rect x="1580" y="14" width="530" height="36" rx="18" fill="#301c48" stroke="#f5c542" stroke-width="1.5" />
    <text x="1845" y="38" text-anchor="middle" font-family="sans-serif" font-size="18" font-weight="bold" fill="#ffe885">✦ 即時詠唱 · 全圖無死角精準轟炸 ✦</text>

    <g transform="translate(40, 58)">
      <rect width="680" height="95" rx="12" fill="#0d0e15" stroke="#f97316" stroke-width="1.5" />
      <image href="{p}scorched.png" xlink:href="{p}scorched.png" x="18" y="15" width="65" height="65" />
      <text x="95" y="42" font-family="serif" font-size="23" font-weight="bold" fill="#ffffff">火焰天罰 (Meteor Fire)</text>
      <text x="95" y="72" font-family="sans-serif" font-size="17" fill="#94a3b8">指定區域天降隕石爆破，高額範圍灼燒！</text>
    </g>
    <g transform="translate(740, 58)">
      <rect width="680" height="95" rx="12" fill="#0d0e15" stroke="#eab308" stroke-width="1.5" />
      <image href="{p}magician.png" xlink:href="{p}magician.png" x="18" y="15" width="65" height="65" />
      <text x="95" y="42" font-family="serif" font-size="23" font-weight="bold" fill="#ffffff">連鎖閃電 (Lightning)</text>
      <text x="95" y="72" font-family="sans-serif" font-size="17" fill="#94a3b8">九天神雷貫穿魔物群，多重彈射瞬秒夜行者！</text>
    </g>
    <g transform="translate(1440, 58)">
      <rect width="680" height="95" rx="12" fill="#0d0e15" stroke="#38bdf8" stroke-width="1.5" />
      <image href="{p}swamp.png" xlink:href="{p}swamp.png" x="18" y="15" width="65" height="65" />
      <text x="95" y="42" font-family="serif" font-size="23" font-weight="bold" fill="#ffffff">極地冰凍 (Freeze)</text>
      <text x="95" y="72" font-family="sans-serif" font-size="17" fill="#94a3b8">3x3 暴風雪大範圍霜結，強制冰封所有魔物！</text>
    </g>

    <line x1="40" y1="170" x2="2120" y2="170" stroke="#f5c542" stroke-width="1" stroke-dasharray="6,6" stroke-opacity="0.4" />
    <text x="40" y="200" font-family="serif" font-size="25" font-weight="bold" fill="#ffe885">🌟 獨特技能點天賦升級系統 (Skill Tree Progression)</text>
    <text x="1560" y="200" font-family="sans-serif" font-size="18" font-weight="bold" fill="#56ccf2">✦ [K] 鍵開啟天賦樹 · 每夜黎明結算 2 SP ✦</text>

    <g transform="translate(40, 215)">
      <rect width="680" height="70" rx="10" fill="#0d0e15" stroke="#68d391" stroke-width="1.5" />
      <image href="{p}farmer.png" xlink:href="{p}farmer.png" x="12" y="10" width="50" height="50" />
      <text x="70" y="35" font-family="serif" font-size="20" font-weight="bold" fill="#a7f3d0">🌾 生產天賦 · 馴獸大師</text>
      <text x="70" y="56" font-family="sans-serif" font-size="15" fill="#cbd5e1">作物速度 +25% · 採集雙倍產出 · 馴化大增</text>
    </g>
    <g transform="translate(740, 215)">
      <rect width="680" height="70" rx="10" fill="#0d0e15" stroke="#f5c542" stroke-width="1.5" />
      <image href="{p}knight.png" xlink:href="{p}knight.png" x="12" y="10" width="50" height="50" />
      <text x="70" y="35" font-family="serif" font-size="20" font-weight="bold" fill="#fef08a">🛡️ 防禦天賦 · 堅甲要塞</text>
      <text x="70" y="56" font-family="sans-serif" font-size="15" fill="#cbd5e1">箭塔射程攻速 +30% · 城牆耐久 · 全員生命</text>
    </g>
    <g transform="translate(1440, 215)">
      <rect width="680" height="70" rx="10" fill="#0d0e15" stroke="#c084fc" stroke-width="1.5" />
      <image href="{p}raw_stone.png" xlink:href="{p}raw_stone.png" x="12" y="10" width="50" height="50" />
      <text x="70" y="35" font-family="serif" font-size="20" font-weight="bold" fill="#f3e8ff">🔮 魔法天賦 · 元素超載</text>
      <text x="70" y="56" font-family="sans-serif" font-size="15" fill="#cbd5e1">全圖冷卻 -25% · 火焰範圍 +50% · 冰凍延長</text>
    </g>
  </g>

  <g id="fence5" transform="translate(100, 2315)">
    <rect x="20" y="45" width="2240" height="8" rx="4" fill="#a06e3b" />
    {fence_svg_row}
  </g>

  <g id="features" transform="translate(0, 2415)">
    <g transform="translate(160, 0)">
      <rect width="1050" height="400" rx="20" fill="#0f111a" stroke="#e63946" stroke-width="2.5" />
      <text x="35" y="48" font-family="serif" font-size="30" font-weight="bold" fill="#ffffff">🌙 暗夜深淵魔物潮與防禦箭塔</text>
      <text x="35" y="85" font-family="sans-serif" font-size="20" fill="#cbd5e1">黑夜降臨巢穴湧出魔物！防禦箭塔全自動遠程齊射，配合全圖元素法術殲滅敵軍！</text>
      <g transform="translate(35, 110)">
        <rect width="225" height="260" rx="14" fill="#181c2b" stroke="#e63946" stroke-width="1.5" />
        <image href="{p}werewolf.png" xlink:href="{p}werewolf.png" x="55" y="20" width="115" height="115" />
        <text x="112" y="175" text-anchor="middle" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">狂暴狼人</text>
        <text x="112" y="210" text-anchor="middle" font-family="sans-serif" font-size="19" font-weight="bold" fill="#e63946">高速突襲</text>
      </g>
      <g transform="translate(285, 110)">
        <rect width="225" height="260" rx="14" fill="#181c2b" stroke="#a855f7" stroke-width="1.5" />
        <image href="{p}vampire.png" xlink:href="{p}vampire.png" x="55" y="20" width="115" height="115" />
        <text x="112" y="175" text-anchor="middle" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">暗夜吸血鬼</text>
        <text x="112" y="210" text-anchor="middle" font-family="sans-serif" font-size="19" font-weight="bold" fill="#a855f7">生命偷取</text>
      </g>
      <g transform="translate(535, 110)">
        <rect width="225" height="260" rx="14" fill="#181c2b" stroke="#888888" stroke-width="1.5" />
        <image href="{p}zombie.png" xlink:href="{p}zombie.png" x="55" y="20" width="115" height="115" />
        <text x="112" y="175" text-anchor="middle" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">劇毒殭屍</text>
        <text x="112" y="210" text-anchor="middle" font-family="sans-serif" font-size="19" font-weight="bold" fill="#888888">攻城肉盾</text>
      </g>
      <g transform="translate(785, 110)">
        <rect width="225" height="260" rx="14" fill="#181c2b" stroke="#f5c542" stroke-width="1.5" />
        <image href="{p}tower.png" xlink:href="{p}tower.png" x="55" y="20" width="115" height="115" />
        <text x="112" y="175" text-anchor="middle" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">防禦箭塔</text>
        <text x="112" y="210" text-anchor="middle" font-family="sans-serif" font-size="19" font-weight="bold" fill="#f5c542">自動齊射</text>
      </g>
    </g>

    <g transform="translate(1270, 0)">
      <rect width="1050" height="400" rx="20" fill="#0f111a" stroke="#2ec4b6" stroke-width="2.5" />
      <text x="35" y="48" font-family="serif" font-size="30" font-weight="bold" fill="#ffffff">🌱 白晝荒野馴化與畜欄經濟</text>
      <text x="35" y="85" font-family="sans-serif" font-size="20" fill="#cbd5e1">探索邊境並馴服野生動物！圈養動物為營地持續產肉，駿馬賦予全體居民 +37.5% 移動加速！</text>
      <g transform="translate(35, 110)">
        <rect width="225" height="260" rx="14" fill="#181c2b" stroke="#2ec4b6" stroke-width="1.5" />
        <image href="{p}boar.png" xlink:href="{p}boar.png" x="55" y="20" width="115" height="115" />
        <text x="112" y="175" text-anchor="middle" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">野豬 (Boar)</text>
        <text x="112" y="210" text-anchor="middle" font-family="sans-serif" font-size="19" font-weight="bold" fill="#2ec4b6">被動產肉</text>
      </g>
      <g transform="translate(285, 110)">
        <rect width="225" height="260" rx="14" fill="#181c2b" stroke="#f5c542" stroke-width="1.5" />
        <image href="{p}horse.png" xlink:href="{p}horse.png" x="55" y="20" width="115" height="115" />
        <text x="112" y="175" text-anchor="middle" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">駿馬 (Horse)</text>
        <text x="112" y="210" text-anchor="middle" font-family="sans-serif" font-size="19" font-weight="bold" fill="#f5c542">全域加速</text>
      </g>
      <g transform="translate(535, 110)">
        <rect width="225" height="260" rx="14" fill="#181c2b" stroke="#56ccf2" stroke-width="1.5" />
        <image href="{p}flying_squirrel.png" xlink:href="{p}flying_squirrel.png" x="55" y="20" width="115" height="115" />
        <text x="112" y="175" text-anchor="middle" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">飛鼠 (Squirrel)</text>
        <text x="112" y="210" text-anchor="middle" font-family="sans-serif" font-size="19" font-weight="bold" fill="#56ccf2">寵物隨行</text>
      </g>
      <g transform="translate(785, 110)">
        <rect width="225" height="260" rx="14" fill="#181c2b" stroke="#e63946" stroke-width="1.5" />
        <image href="{p}bear.png" xlink:href="{p}bear.png" x="55" y="20" width="115" height="115" />
        <text x="112" y="175" text-anchor="middle" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">巨熊 (Bear)</text>
        <text x="112" y="210" text-anchor="middle" font-family="sans-serif" font-size="19" font-weight="bold" fill="#e63946">深山霸主</text>
      </g>
    </g>
  </g>

  <g id="fence6" transform="translate(100, 2835)">
    <rect x="20" y="45" width="2240" height="8" rx="4" fill="#a06e3b" />
    {fence_svg_row}
  </g>

  <g id="biomes" transform="translate(160, 2935)">
    <rect width="2160" height="420" rx="20" fill="#0f111a" stroke="#f5c542" stroke-width="2" />
    <text x="40" y="45" font-family="serif" font-size="30" font-weight="bold" fill="#ffffff">🗺️ 16-Tile 四大多元地貌與五大核心建築</text>
    <g transform="translate(40, 70)">
      <g transform="translate(0, 0)">
        <rect width="495" height="150" rx="14" fill="#181c2b" stroke="#56ccf2" stroke-width="1.5" />
        <image href="{p}river.png" xlink:href="{p}river.png" x="25" y="20" width="60" height="60" />
        <text x="105" y="55" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">水彩河流</text>
        <text x="105" y="85" font-family="sans-serif" font-size="18" font-weight="bold" fill="#56ccf2">自然緩速 50%</text>
        <text x="25" y="125" font-family="sans-serif" font-size="16" fill="#cbd5e1">16-Tile 無縫水彩拼接，阻礙敵軍！</text>
      </g>
      <g transform="translate(535, 0)">
        <rect width="495" height="150" rx="14" fill="#181c2b" stroke="#b8860b" stroke-width="1.5" />
        <image href="{p}mountain.png" xlink:href="{p}mountain.png" x="25" y="20" width="60" height="60" />
        <text x="105" y="55" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">連綿山脈</text>
        <text x="105" y="85" font-family="sans-serif" font-size="18" font-weight="bold" fill="#b8860b">泛洪全揭露障礙</text>
        <text x="25" y="125" font-family="sans-serif" font-size="16" fill="#cbd5e1">山腳探索全揭露，山體純淨壯闊！</text>
      </g>
      <g transform="translate(1070, 0)">
        <rect width="495" height="150" rx="14" fill="#181c2b" stroke="#2ec4b6" stroke-width="1.5" />
        <image href="{p}swamp.png" xlink:href="{p}swamp.png" x="25" y="20" width="60" height="60" />
        <text x="105" y="55" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">幽暗泥沼</text>
        <text x="105" y="85" font-family="sans-serif" font-size="18" font-weight="bold" fill="#2ec4b6">5.0s 受困倒數</text>
        <text x="25" y="125" font-family="sans-serif" font-size="16" fill="#cbd5e1">泥濘陷阱限制機動，降低調度！</text>
      </g>
      <g transform="translate(1605, 0)">
        <rect width="495" height="150" rx="14" fill="#181c2b" stroke="#ff7b00" stroke-width="1.5" />
        <image href="{p}scorched.png" xlink:href="{p}scorched.png" x="25" y="20" width="60" height="60" />
        <text x="105" y="55" font-family="serif" font-size="25" font-weight="bold" fill="#ffffff">熔岩焦土</text>
        <text x="105" y="85" font-family="sans-serif" font-size="18" font-weight="bold" fill="#ff7b00">高溫燃燒灼燒</text>
        <text x="25" y="125" font-family="sans-serif" font-size="16" fill="#cbd5e1">地熱熔岩持續扣血，高溫餘燼！</text>
      </g>
    </g>

    <g transform="translate(40, 240)">
      <rect width="2080" height="155" rx="14" fill="#141824" stroke="#f5c542" stroke-width="1.5" />
      <g transform="translate(30, 15)">
        <g transform="translate(0, 0)">
          <rect width="480" height="125" rx="10" fill="#1e2336" stroke="#a0a6b8" stroke-width="1" />
          <image href="{p}wall.png" xlink:href="{p}wall.png" x="15" y="15" width="45" height="45" />
          <text x="75" y="42" font-family="serif" font-size="22" font-weight="bold" fill="#ffffff">城牆 (Wall)</text>
          <text x="75" y="68" font-family="sans-serif" font-size="15" fill="#f5c542">🪵 4 木材 · 阻擋路線</text>
          <text x="15" y="105" font-family="sans-serif" font-size="15" fill="#cbd5e1">保護核心設施，引導魔物！</text>
        </g>
        <g transform="translate(515, 0)">
          <rect width="480" height="125" rx="10" fill="#1e2336" stroke="#f5c542" stroke-width="1" />
          <image href="{p}tower.png" xlink:href="{p}tower.png" x="15" y="15" width="45" height="45" />
          <text x="75" y="42" font-family="serif" font-size="22" font-weight="bold" fill="#ffffff">防禦塔 (Tower)</text>
          <text x="75" y="68" font-family="sans-serif" font-size="15" fill="#f5c542">🪵 2 木 + 🧱 3 磚</text>
          <text x="15" y="105" font-family="sans-serif" font-size="15" fill="#cbd5e1">全自動遠程齊射，夜間支柱！</text>
        </g>
        <g transform="translate(1030, 0)">
          <rect width="480" height="125" rx="10" fill="#1e2336" stroke="#56ccf2" stroke-width="1" />
          <image href="{p}house.png" xlink:href="{p}house.png" x="15" y="15" width="45" height="45" />
          <text x="75" y="42" font-family="serif" font-size="22" font-weight="bold" fill="#ffffff">民房 (House)</text>
          <text x="75" y="68" font-family="sans-serif" font-size="15" fill="#f5c542">🪵 4 木 + 🧱 2 磚</text>
          <text x="15" y="105" font-family="sans-serif" font-size="15" fill="#cbd5e1">+1 人口，黎明誕生拓荒者！</text>
        </g>
        <g transform="translate(1545, 0)">
          <rect width="480" height="125" rx="10" fill="#1e2336" stroke="#2ec4b6" stroke-width="1" />
          <image href="{p}farmland.png" xlink:href="{p}farmland.png" x="15" y="15" width="45" height="45" />
          <text x="75" y="42" font-family="serif" font-size="22" font-weight="bold" fill="#ffffff">農田 (Farmland)</text>
          <text x="75" y="68" font-family="sans-serif" font-size="15" fill="#f5c542">🪵 2 木 + 🌾 1 作物</text>
          <text x="15" y="105" font-family="sans-serif" font-size="15" fill="#cbd5e1">成熟收割，永續糧食命脈！</text>
        </g>
      </g>
    </g>
  </g>

  <g id="footer" transform="translate(1240, 3420)" text-anchor="middle">
    <rect x="-600" y="-30" width="1200" height="60" rx="16" fill="#101320" stroke="#f5c542" stroke-width="1" />
    <text x="0" y="10" font-family="serif" font-size="30" font-weight="bold" fill="#f5c542">第二組 出品 · 《MEDIEVIL》</text>
  </g>
</svg>"""

with open('docs/a1_poster_editable.svg', 'w', encoding='utf-8') as f:
    f.write(get_svg_content('assets/'))
with open('a1_poster_editable.svg', 'w', encoding='utf-8') as f:
    f.write(get_svg_content('docs/assets/'))
print("[OK] Successfully generated lightweight a1_poster_editable.svg (~35 KB)")

# ── 4. NATIVE EDITABLE POWERPOINT PRESENTATION (.pptx) ───────────────────────
prs = Presentation()
prs.slide_width = Inches(23.38)
prs.slide_height = Inches(33.11)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(23.38), Inches(33.11))
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(8, 9, 16)
bg.line.color.rgb = RGBColor(245, 197, 66)
bg.line.width = Pt(4)

title_box = slide.shapes.add_textbox(Inches(2), Inches(0.8), Inches(19.38), Inches(2.2))
tf = title_box.text_frame
tf.word_wrap = True
p0 = tf.paragraphs[0]
p0.text = "✦ 2026 年度暗黑奇幻像素殖民地生存策略 ✦"
p0.font.size = Pt(22)
p0.font.bold = True
p0.font.color.rgb = RGBColor(255, 232, 133)
p0.alignment = PP_ALIGN.CENTER

p1 = tf.add_paragraph()
p1.text = "M E D I E V I L"
p1.font.size = Pt(66)
p1.font.bold = True
p1.font.color.rgb = RGBColor(245, 197, 66)
p1.alignment = PP_ALIGN.CENTER

parade_items = [
    ("docs/assets/farmer.png", "拓荒農夫", Inches(1.8)),
    ("docs/assets/horse.png", "疾行駿馬", Inches(3.8)),
    ("docs/assets/boar.png", "野豬產肉", Inches(5.8)),
    ("docs/assets/knight.png", "聖殿騎士", Inches(7.8)),
    ("docs/assets/magician.png", "元素法師", Inches(9.8)),
    ("docs/assets/flying_squirrel.png", "隨行飛鼠", Inches(11.8)),
    ("docs/assets/werewolf.png", "狂暴狼人", Inches(13.8)),
    ("docs/assets/vampire.png", "暗夜吸血鬼", Inches(15.8)),
    ("docs/assets/zombie.png", "劇毒殭屍", Inches(17.8)),
    ("docs/assets/bear.png", "狂暴巨熊", Inches(19.8)),
]

for img_p, lbl, x_pos in parade_items:
    if os.path.exists(img_p):
        slide.shapes.add_picture(img_p, x_pos, Inches(2.8), width=Inches(1.3))
    p_box = slide.shapes.add_textbox(x_pos - Inches(0.2), Inches(4.2), Inches(1.7), Inches(0.6))
    p_box.text_frame.paragraphs[0].text = lbl
    p_box.text_frame.paragraphs[0].font.size = Pt(14)
    p_box.text_frame.paragraphs[0].font.bold = True
    p_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)
    p_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

lore_box = slide.shapes.add_textbox(Inches(2), Inches(4.9), Inches(19.38), Inches(1.0))
ltf = lore_box.text_frame
ltf.word_wrap = True
lp = ltf.paragraphs[0]
lp.text = "「白晝在未知的迷霧荒野中拓荒建設、採集資源、耕種農田、馴服野獸；黑夜深淵魔物破曉襲擊之時，詠唱全圖魔法抵禦狂暴狼人、吸血鬼與劇毒殭屍的無盡圍攻。白晝拓荒寸土，黑夜寸步不讓！」"
lp.font.size = Pt(17)
lp.font.color.rgb = RGBColor(197, 203, 216)
lp.alignment = PP_ALIGN.CENTER

if os.path.exists("docs/assets/screenshot_day.png"):
    slide.shapes.add_picture("docs/assets/screenshot_day.png", Inches(1.8), Inches(6.0), width=Inches(9.5))
if os.path.exists("docs/assets/screenshot_night.png"):
    slide.shapes.add_picture("docs/assets/screenshot_night.png", Inches(12.0), Inches(6.0), width=Inches(9.5))

task_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(14.5), Inches(19.78), Inches(2.2))
task_card.fill.solid()
task_card.fill.fore_color.rgb = RGBColor(16, 37, 32)
task_card.line.color.rgb = RGBColor(46, 196, 182)
task_card.line.width = Pt(2.5)

ttf = task_card.text_frame
ttf.word_wrap = True
tp0 = ttf.paragraphs[0]
tp0.text = "📋⚡【高度自由的任務安排優先順序，由你決定！！】"
tp0.font.size = Pt(22)
tp0.font.bold = True
tp0.font.color.rgb = RGBColor(167, 243, 208)

tp1 = ttf.add_paragraph()
tp1.text = "🌾 [農耕 vs 伐木採石] 自由設定居民工作優先級！  |  🐗 [狩獵 vs 馴化畜牧] 優先獵肉或圈養加速？  |  🏰 [要塞 vs 巡邏] 白晝搶修防禦塔，入夜騎士阻擊！"
tp1.font.size = Pt(15.5)
tp1.font.color.rgb = RGBColor(203, 213, 225)

magic_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(17.0), Inches(19.78), Inches(3.2))
magic_card.fill.solid()
magic_card.fill.fore_color.rgb = RGBColor(26, 17, 40)
magic_card.line.color.rgb = RGBColor(168, 85, 247)
magic_card.line.width = Pt(2.5)

mtf = magic_card.text_frame
mtf.word_wrap = True
mp0 = mtf.paragraphs[0]
mp0.text = "⚡🔥❄️ 全圖三大元素法術天罰打擊 (Global Magic Strike)"
mp0.font.size = Pt(22)
mp0.font.bold = True
mp0.font.color.rgb = RGBColor(243, 232, 255)

mp1 = mtf.add_paragraph()
mp1.text = "🔥 [火焰天罰] 隕石轟炸爆破！  |  ⚡ [連鎖閃電] 九天神雷多重彈射！  |  ❄️ [極地冰凍] 3x3 霜結強制冰封！"
mp1.font.size = Pt(15.5)
mp1.font.color.rgb = RGBColor(203, 213, 225)

mp2 = mtf.add_paragraph()
mp2.text = "🌟 獨特技能點天賦升級系統 (Skill Tree Progression) ── 按 [K] 鍵開啟天賦樹 · 每夜黎明結算 2 SP 點數"
mp2.font.size = Pt(20)
mp2.font.bold = True
mp2.font.color.rgb = RGBColor(255, 232, 133)

mp3 = mtf.add_paragraph()
mp3.text = "🌾 [生產天賦] 作物生長 +25% · 雙倍採集 · 馴獸大師  |  🛡️ [防禦天賦] 箭塔射程攻速 +30% · 城牆耐久 · 全員生命  |  🔮 [魔法天賦] 冷卻 -25% · 火焰範圍 +50% · 冰凍延長"
mp3.font.size = Pt(15.5)
mp3.font.color.rgb = RGBColor(203, 213, 225)

footer_box = slide.shapes.add_textbox(Inches(2), Inches(31.5), Inches(19.38), Inches(0.8))
ftf = footer_box.text_frame
ftf.word_wrap = True
fp = ftf.paragraphs[0]
fp.text = "第二組 出品 · 《MEDIEVIL》"
fp.font.size = Pt(18)
fp.font.bold = True
fp.font.color.rgb = RGBColor(245, 197, 66)
fp.alignment = PP_ALIGN.CENTER

prs.save("a1_poster_editable.pptx")
prs.save("docs/a1_poster_editable.pptx")
print(f"[OK] Successfully generated a1_poster_editable.pptx ({os.path.getsize('a1_poster_editable.pptx')//1024} KB)")

print("All Lightweight Files Generated Successfully!")
